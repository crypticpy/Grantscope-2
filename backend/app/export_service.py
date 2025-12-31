"""
Export Service for Intelligence Cards and Workstream Reports.

This service handles generation of export files in multiple formats:
- PDF: Using ReportLab for professional document generation
- PowerPoint: Using python-pptx for presentation slides
- CSV: Using pandas for tabular data export

Chart Generation:
- Score radar charts showing all dimensions
- Score bar charts for individual scores
- Pillar distribution charts for workstream reports
- All charts use matplotlib with 'Agg' backend (non-GUI)

Usage:
    export_service = ExportService(supabase_client)
    pdf_path = await export_service.generate_pdf(card_data)
    pptx_path = await export_service.generate_pptx(card_data)
    csv_content = await export_service.generate_csv(card_data)
"""

import io
import logging
import os
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple

import matplotlib
matplotlib.use('Agg')  # Non-GUI backend - must be set before importing pyplot
import matplotlib.pyplot as plt
import numpy as np

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ReportLab imports for PDF generation
from reportlab.lib import colors as rl_colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    Image as RLImage,
    PageBreak,
    HRFlowable,
    BaseDocTemplate,
    Frame,
    PageTemplate,
    KeepTogether,
)
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
from reportlab.pdfgen import canvas

from supabase import Client

from .models.export import (
    CardExportData,
    ExportFormat,
    EXPORT_CONTENT_TYPES,
    get_export_filename,
)

logger = logging.getLogger(__name__)


# ============================================================================
# Constants
# ============================================================================

# Foresight branding colors
FORESIGHT_COLORS = {
    "primary": "#1E3A5F",      # Deep blue
    "secondary": "#2E86AB",    # Lighter blue
    "accent": "#A23B72",       # Magenta accent
    "success": "#28A745",      # Green for positive metrics
    "warning": "#FFC107",      # Yellow for warnings
    "danger": "#DC3545",       # Red for negative/risk
    "light": "#F8F9FA",        # Light background
    "dark": "#343A40",         # Dark text
}

# Score colors for charts
SCORE_COLORS = {
    "Novelty": "#2E86AB",
    "Maturity": "#28A745",
    "Impact": "#A23B72",
    "Relevance": "#FFC107",
    "Velocity": "#17A2B8",
    "Risk": "#DC3545",
    "Opportunity": "#6F42C1",
}

# Chart settings
CHART_DPI = 300
CHART_FIGURE_SIZE = (8, 6)
RADAR_FIGURE_SIZE = (8, 8)

# PowerPoint settings
PPTX_SLIDE_WIDTH = Inches(13.333)  # 16:9 widescreen
PPTX_SLIDE_HEIGHT = Inches(7.5)
PPTX_TITLE_FONT_SIZE = Pt(44)
PPTX_SUBTITLE_FONT_SIZE = Pt(24)
PPTX_BODY_FONT_SIZE = Pt(18)
PPTX_SMALL_FONT_SIZE = Pt(14)
PPTX_MARGIN = Inches(0.5)
PPTX_CHART_WIDTH = Inches(5)
PPTX_CHART_HEIGHT = Inches(4)

# PDF settings
PDF_PAGE_SIZE = letter
PDF_MARGIN = 0.75 * inch
PDF_TITLE_FONT_SIZE = 24
PDF_HEADING_FONT_SIZE = 14
PDF_BODY_FONT_SIZE = 11
PDF_SMALL_FONT_SIZE = 9
PDF_CHART_WIDTH = 5.5 * inch
PDF_CHART_HEIGHT = 4 * inch

# ReportLab color conversion helper
def hex_to_rl_color(hex_color: str) -> rl_colors.Color:
    """Convert hex color string to ReportLab Color object."""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16) / 255.0
    g = int(hex_color[2:4], 16) / 255.0
    b = int(hex_color[4:6], 16) / 255.0
    return rl_colors.Color(r, g, b)


# PDF color palette using Foresight colors
PDF_COLORS = {
    "primary": hex_to_rl_color(FORESIGHT_COLORS["primary"]),
    "secondary": hex_to_rl_color(FORESIGHT_COLORS["secondary"]),
    "accent": hex_to_rl_color(FORESIGHT_COLORS["accent"]),
    "success": hex_to_rl_color(FORESIGHT_COLORS["success"]),
    "warning": hex_to_rl_color(FORESIGHT_COLORS["warning"]),
    "danger": hex_to_rl_color(FORESIGHT_COLORS["danger"]),
    "light": hex_to_rl_color(FORESIGHT_COLORS["light"]),
    "dark": hex_to_rl_color(FORESIGHT_COLORS["dark"]),
}

# ============================================================================
# Branding Assets & AI Disclosure
# ============================================================================

# Path to City of Austin logo (relative to backend directory)
def _get_logo_path() -> Optional[str]:
    """Get the path to the City of Austin logo, checking multiple locations."""
    # Check relative to this file's location
    base_dir = Path(__file__).parent.parent.parent  # Go up from app/ to backend/ to project root
    possible_paths = [
        base_dir / "branding" / "COA-Logo-Horizontal-Official-RGB.png",
        base_dir / "branding" / "COA-Icon-Official-RGB.png",
        Path("/app/branding/COA-Logo-Horizontal-Official-RGB.png"),  # Railway deployment
        Path("/app/branding/COA-Icon-Official-RGB.png"),  # Railway deployment fallback
    ]
    
    for path in possible_paths:
        if path.exists():
            return str(path)
    
    logger.warning("City of Austin logo not found in expected locations")
    return None

COA_LOGO_PATH = _get_logo_path()

# AI Technology Disclosure - transparent listing of all AI systems used
AI_TECHNOLOGY_DISCLOSURE = """
AI-Powered Research & Analysis Platform

This strategic intelligence report was generated using advanced artificial intelligence technologies:

• Anthropic Claude (Opus/Sonnet) - Strategic analysis, synthesis, and report generation
• OpenAI GPT-4o - Classification, scoring, and natural language processing  
• GPT Researcher - Autonomous deep research and source discovery
• Exa AI - High-quality source search and content retrieval
• Firecrawl - Web content extraction and PDF processing
• Tavily - Real-time news and research aggregation

The City of Austin is committed to transparent and responsible use of AI technology 
in public service. All AI-generated content is reviewed for accuracy and relevance 
to municipal government operations.
""".strip()

# Shorter disclosure for footer
AI_DISCLOSURE_SHORT = (
    "Generated by Foresight Strategic Intelligence Platform | "
    "AI Technologies: Anthropic Claude, OpenAI GPT-4o, GPT Researcher, Exa AI, Firecrawl, Tavily"
)


# ============================================================================
# Professional PDF Builder with Header/Footer
# ============================================================================

class ProfessionalPDFBuilder:
    """
    Custom PDF document builder with professional header and footer.
    
    Used for executive briefs and workstream reports that will be shared
    with senior city leadership.
    """
    
    def __init__(
        self,
        filename: str,
        title: str,
        subtitle: Optional[str] = None,
        include_logo: bool = True,
        include_ai_disclosure: bool = True
    ):
        """
        Initialize the PDF builder.
        
        Args:
            filename: Output PDF file path
            title: Document title for header
            subtitle: Optional subtitle
            include_logo: Whether to include City of Austin logo
            include_ai_disclosure: Whether to include AI technology disclosure
        """
        self.filename = filename
        self.title = title
        self.subtitle = subtitle
        self.include_logo = include_logo
        self.include_ai_disclosure = include_ai_disclosure
        self.page_width, self.page_height = letter
        
        # Margins - extra space at top/bottom for header/footer
        self.left_margin = 0.75 * inch
        self.right_margin = 0.75 * inch
        self.top_margin = 1.4 * inch  # Extra space for header
        self.bottom_margin = 1.2 * inch  # Extra space for footer
        
    def _draw_header(self, canvas_obj: canvas.Canvas, doc):
        """Draw the professional header on each page."""
        canvas_obj.saveState()
        
        # Header background - subtle gradient effect using rectangle
        canvas_obj.setFillColor(PDF_COLORS["primary"])
        canvas_obj.rect(0, self.page_height - 1.1 * inch, self.page_width, 1.1 * inch, fill=True, stroke=False)
        
        # Add accent line below header
        canvas_obj.setStrokeColor(hex_to_rl_color("#2E86AB"))
        canvas_obj.setLineWidth(3)
        canvas_obj.line(0, self.page_height - 1.1 * inch, self.page_width, self.page_height - 1.1 * inch)
        
        # City of Austin logo (if available and enabled)
        logo_width = 0
        if self.include_logo and COA_LOGO_PATH and Path(COA_LOGO_PATH).exists():
            try:
                # Draw logo on left side of header
                logo_height = 0.6 * inch
                logo_width = 1.8 * inch  # Aspect ratio ~3:1 for horizontal logo
                canvas_obj.drawImage(
                    COA_LOGO_PATH,
                    self.left_margin,
                    self.page_height - 0.9 * inch,
                    width=logo_width,
                    height=logo_height,
                    preserveAspectRatio=True,
                    mask='auto'
                )
                logo_width += 0.3 * inch  # Add spacing after logo
            except Exception as e:
                logger.warning(f"Failed to add logo to PDF header: {e}")
                logo_width = 0
        
        # Foresight branding text
        text_x = self.left_margin + logo_width
        
        # "FORESIGHT" title
        canvas_obj.setFillColor(rl_colors.white)
        canvas_obj.setFont("Helvetica-Bold", 16)
        canvas_obj.drawString(text_x, self.page_height - 0.45 * inch, "FORESIGHT")
        
        # "Strategic Intelligence" subtitle
        canvas_obj.setFont("Helvetica", 10)
        canvas_obj.drawString(text_x, self.page_height - 0.65 * inch, "Strategic Intelligence Platform")
        
        # Document title on right side
        canvas_obj.setFont("Helvetica-Bold", 11)
        title_text = self.title[:50] + "..." if len(self.title) > 50 else self.title
        title_width = canvas_obj.stringWidth(title_text, "Helvetica-Bold", 11)
        canvas_obj.drawString(
            self.page_width - self.right_margin - title_width,
            self.page_height - 0.45 * inch,
            title_text
        )
        
        # Page generation date on right
        canvas_obj.setFont("Helvetica", 9)
        date_text = datetime.now().strftime("%B %d, %Y")
        date_width = canvas_obj.stringWidth(date_text, "Helvetica", 9)
        canvas_obj.drawString(
            self.page_width - self.right_margin - date_width,
            self.page_height - 0.65 * inch,
            date_text
        )
        
        canvas_obj.restoreState()
    
    def _draw_footer(self, canvas_obj: canvas.Canvas, doc):
        """Draw the professional footer on each page."""
        canvas_obj.saveState()
        
        # Footer background
        canvas_obj.setFillColor(hex_to_rl_color("#F8F9FA"))
        canvas_obj.rect(0, 0, self.page_width, 1.0 * inch, fill=True, stroke=False)
        
        # Top border line
        canvas_obj.setStrokeColor(PDF_COLORS["primary"])
        canvas_obj.setLineWidth(1)
        canvas_obj.line(self.left_margin, 1.0 * inch, self.page_width - self.right_margin, 1.0 * inch)
        
        # AI Technology Disclosure (if enabled)
        if self.include_ai_disclosure:
            canvas_obj.setFillColor(PDF_COLORS["dark"])
            canvas_obj.setFont("Helvetica", 7)
            
            # Multi-line disclosure
            disclosure_lines = [
                "AI-Powered Intelligence: This report was generated using Anthropic Claude, OpenAI GPT-4o,",
                "GPT Researcher, Exa AI, Firecrawl, and Tavily. The City of Austin is committed to",
                "transparent and responsible use of AI technology in public service."
            ]
            
            y_pos = 0.75 * inch
            for line in disclosure_lines:
                canvas_obj.drawString(self.left_margin, y_pos, line)
                y_pos -= 0.12 * inch
        
        # Page number on right
        canvas_obj.setFont("Helvetica", 9)
        page_text = f"Page {doc.page}"
        page_width = canvas_obj.stringWidth(page_text, "Helvetica", 9)
        canvas_obj.drawString(
            self.page_width - self.right_margin - page_width,
            0.4 * inch,
            page_text
        )
        
        # Confidentiality notice
        canvas_obj.setFont("Helvetica-Oblique", 7)
        canvas_obj.setFillColor(rl_colors.gray)
        canvas_obj.drawString(
            self.left_margin,
            0.25 * inch,
            "City of Austin Internal Document - For Official Use"
        )
        
        canvas_obj.restoreState()
    
    def _draw_header_footer(self, canvas_obj: canvas.Canvas, doc):
        """Combined callback for header and footer."""
        self._draw_header(canvas_obj, doc)
        self._draw_footer(canvas_obj, doc)
    
    def build(self, elements: List[Any]) -> str:
        """
        Build the PDF document with all elements.
        
        Args:
            elements: List of ReportLab flowable elements
            
        Returns:
            Path to the generated PDF file
        """
        # Create document with custom page template
        doc = BaseDocTemplate(
            self.filename,
            pagesize=letter,
            leftMargin=self.left_margin,
            rightMargin=self.right_margin,
            topMargin=self.top_margin,
            bottomMargin=self.bottom_margin
        )
        
        # Create frame for content
        frame = Frame(
            self.left_margin,
            self.bottom_margin,
            self.page_width - self.left_margin - self.right_margin,
            self.page_height - self.top_margin - self.bottom_margin,
            id='normal'
        )
        
        # Create page template with header/footer
        template = PageTemplate(
            id='professional',
            frames=[frame],
            onPage=self._draw_header_footer
        )
        
        doc.addPageTemplates([template])
        doc.build(elements)
        
        return self.filename


def get_professional_pdf_styles() -> Dict[str, ParagraphStyle]:
    """
    Create custom paragraph styles for professional PDF generation.
    
    Returns:
        Dictionary of ParagraphStyle objects for executive documents
    """
    styles = getSampleStyleSheet()
    
    return {
        'DocTitle': ParagraphStyle(
            'DocTitle',
            parent=styles['Heading1'],
            fontSize=28,
            textColor=PDF_COLORS["primary"],
            spaceAfter=16,
            spaceBefore=8,
            alignment=TA_LEFT,
            fontName='Helvetica-Bold',
            leading=34,
        ),
        'DocSubtitle': ParagraphStyle(
            'DocSubtitle',
            parent=styles['Normal'],
            fontSize=14,
            textColor=PDF_COLORS["secondary"],
            spaceAfter=20,
            fontName='Helvetica',
            leading=18,
        ),
        'SectionHeading': ParagraphStyle(
            'SectionHeading',
            parent=styles['Heading1'],
            fontSize=16,
            textColor=PDF_COLORS["primary"],
            spaceBefore=20,
            spaceAfter=10,
            fontName='Helvetica-Bold',
            borderPadding=(0, 0, 4, 0),
        ),
        'SubsectionHeading': ParagraphStyle(
            'SubsectionHeading',
            parent=styles['Heading2'],
            fontSize=13,
            textColor=PDF_COLORS["secondary"],
            spaceBefore=14,
            spaceAfter=6,
            fontName='Helvetica-Bold',
        ),
        'BodyText': ParagraphStyle(
            'BodyText',
            parent=styles['Normal'],
            fontSize=11,
            textColor=PDF_COLORS["dark"],
            spaceBefore=4,
            spaceAfter=8,
            leading=15,
            alignment=TA_JUSTIFY,
        ),
        'BulletText': ParagraphStyle(
            'BulletText',
            parent=styles['Normal'],
            fontSize=11,
            textColor=PDF_COLORS["dark"],
            spaceBefore=2,
            spaceAfter=4,
            leftIndent=20,
            bulletIndent=10,
            leading=14,
        ),
        'SmallText': ParagraphStyle(
            'SmallText',
            parent=styles['Normal'],
            fontSize=9,
            textColor=rl_colors.gray,
            spaceBefore=2,
            spaceAfter=2,
        ),
        'MetadataText': ParagraphStyle(
            'MetadataText',
            parent=styles['Normal'],
            fontSize=10,
            textColor=PDF_COLORS["dark"],
            spaceBefore=2,
            spaceAfter=2,
            fontName='Helvetica',
        ),
        'CalloutText': ParagraphStyle(
            'CalloutText',
            parent=styles['Normal'],
            fontSize=11,
            textColor=PDF_COLORS["primary"],
            spaceBefore=8,
            spaceAfter=8,
            leftIndent=15,
            rightIndent=15,
            borderPadding=10,
            backColor=hex_to_rl_color("#F0F4F8"),
            leading=15,
        ),
        'ExecutiveSummary': ParagraphStyle(
            'ExecutiveSummary',
            parent=styles['Normal'],
            fontSize=12,
            textColor=PDF_COLORS["dark"],
            spaceBefore=6,
            spaceAfter=12,
            leading=17,
            alignment=TA_JUSTIFY,
            fontName='Helvetica',
        ),
    }


# ============================================================================
# Export Service
# ============================================================================

class ExportService:
    """
    Service for generating export files from intelligence cards and workstreams.

    Supports PDF, PowerPoint, and CSV export formats with embedded visualizations.
    Follows the service class pattern from research_service.py.
    """

    def __init__(self, supabase: Client):
        """
        Initialize the ExportService.

        Args:
            supabase: Supabase client for database queries
        """
        self.supabase = supabase
        logger.info("ExportService initialized")

    # ========================================================================
    # Chart Generation Methods
    # ========================================================================

    def generate_score_chart(
        self,
        card_data: CardExportData,
        chart_type: str = "bar",
        dpi: int = CHART_DPI
    ) -> Optional[str]:
        """
        Generate a chart showing card scores.

        Args:
            card_data: Card data containing scores
            chart_type: Type of chart ('bar' or 'radar')
            dpi: Resolution for the chart image

        Returns:
            Path to the generated chart image, or None if generation fails
        """
        try:
            scores = card_data.get_all_scores()

            # Filter out None scores
            valid_scores = {k: v for k, v in scores.items() if v is not None}

            if not valid_scores:
                logger.warning(f"No valid scores for card {card_data.id}, skipping chart")
                return None

            if chart_type == "radar":
                return self._generate_radar_chart(valid_scores, card_data.name, dpi)
            else:
                return self._generate_bar_chart(valid_scores, card_data.name, dpi)

        except Exception as e:
            logger.error(f"Error generating score chart: {e}")
            return None

    def _generate_bar_chart(
        self,
        scores: Dict[str, int],
        title: str,
        dpi: int
    ) -> str:
        """
        Generate a horizontal bar chart of scores.

        Args:
            scores: Dictionary of score names to values
            title: Chart title
            dpi: Resolution for the image

        Returns:
            Path to the generated chart image
        """
        fig, ax = plt.subplots(figsize=CHART_FIGURE_SIZE)

        try:
            labels = list(scores.keys())
            values = list(scores.values())
            colors = [SCORE_COLORS.get(label, FORESIGHT_COLORS["primary"]) for label in labels]

            y_pos = np.arange(len(labels))

            bars = ax.barh(y_pos, values, color=colors, edgecolor='white', height=0.6)

            # Add value labels on bars
            for bar, value in zip(bars, values):
                width = bar.get_width()
                ax.text(
                    width + 2, bar.get_y() + bar.get_height() / 2,
                    f'{value}',
                    va='center', ha='left',
                    fontsize=10, fontweight='bold',
                    color=FORESIGHT_COLORS["dark"]
                )

            ax.set_yticks(y_pos)
            ax.set_yticklabels(labels, fontsize=11)
            ax.set_xlim(0, 110)  # Extra space for labels
            ax.set_xlabel('Score (0-100)', fontsize=11)
            ax.set_title(f'Scores: {title[:40]}...', fontsize=12, fontweight='bold', pad=15)

            # Style the chart
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['bottom'].set_color(FORESIGHT_COLORS["light"])
            ax.spines['left'].set_color(FORESIGHT_COLORS["light"])

            # Add gridlines
            ax.xaxis.grid(True, linestyle='--', alpha=0.3)
            ax.set_axisbelow(True)

            plt.tight_layout()

            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix='.png',
                delete=False,
                prefix='foresight_chart_'
            )
            plt.savefig(temp_file.name, dpi=dpi, bbox_inches='tight', facecolor='white')

            return temp_file.name

        finally:
            plt.close(fig)  # CRITICAL: Prevent memory leaks

    def _generate_radar_chart(
        self,
        scores: Dict[str, int],
        title: str,
        dpi: int
    ) -> str:
        """
        Generate a radar/spider chart of scores.

        Args:
            scores: Dictionary of score names to values
            title: Chart title
            dpi: Resolution for the image

        Returns:
            Path to the generated chart image
        """
        fig, ax = plt.subplots(figsize=RADAR_FIGURE_SIZE, subplot_kw=dict(polar=True))

        try:
            labels = list(scores.keys())
            values = list(scores.values())

            # Number of variables
            num_vars = len(labels)

            # Compute angle for each axis
            angles = [n / float(num_vars) * 2 * np.pi for n in range(num_vars)]
            angles += angles[:1]  # Complete the loop

            # Complete the data loop
            values_plot = values + values[:1]

            # Plot the data
            ax.plot(angles, values_plot, 'o-', linewidth=2, color=FORESIGHT_COLORS["primary"])
            ax.fill(angles, values_plot, alpha=0.25, color=FORESIGHT_COLORS["secondary"])

            # Set the labels
            ax.set_xticks(angles[:-1])
            ax.set_xticklabels(labels, fontsize=11)

            # Set y-axis limits
            ax.set_ylim(0, 100)
            ax.set_yticks([20, 40, 60, 80, 100])
            ax.set_yticklabels(['20', '40', '60', '80', '100'], fontsize=9, color='gray')

            # Add title
            ax.set_title(
                f'Score Profile: {title[:35]}...',
                fontsize=12, fontweight='bold', pad=20
            )

            plt.tight_layout()

            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix='.png',
                delete=False,
                prefix='foresight_radar_'
            )
            plt.savefig(temp_file.name, dpi=dpi, bbox_inches='tight', facecolor='white')

            return temp_file.name

        finally:
            plt.close(fig)  # CRITICAL: Prevent memory leaks

    def generate_pillar_distribution_chart(
        self,
        pillar_counts: Dict[str, int],
        title: str = "Pillar Distribution",
        dpi: int = CHART_DPI
    ) -> Optional[str]:
        """
        Generate a pie/donut chart showing distribution of cards across pillars.

        Args:
            pillar_counts: Dictionary mapping pillar names to card counts
            title: Chart title
            dpi: Resolution for the image

        Returns:
            Path to the generated chart image, or None if no data
        """
        if not pillar_counts:
            logger.warning("No pillar data for distribution chart")
            return None

        fig, ax = plt.subplots(figsize=CHART_FIGURE_SIZE)

        try:
            labels = list(pillar_counts.keys())
            values = list(pillar_counts.values())

            # Generate colors from palette
            colors = plt.cm.Set2(np.linspace(0, 1, len(labels)))

            # Create donut chart
            wedges, texts, autotexts = ax.pie(
                values,
                labels=labels,
                autopct='%1.1f%%',
                colors=colors,
                pctdistance=0.75,
                wedgeprops=dict(width=0.5, edgecolor='white'),
                textprops={'fontsize': 10}
            )

            # Style the percentage text
            for autotext in autotexts:
                autotext.set_fontsize(9)
                autotext.set_fontweight('bold')

            ax.set_title(title, fontsize=12, fontweight='bold', pad=15)

            # Add legend
            ax.legend(
                wedges, [f'{l} ({v})' for l, v in zip(labels, values)],
                title="Pillars",
                loc="center left",
                bbox_to_anchor=(1, 0, 0.5, 1),
                fontsize=9
            )

            plt.tight_layout()

            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix='.png',
                delete=False,
                prefix='foresight_pillar_'
            )
            plt.savefig(temp_file.name, dpi=dpi, bbox_inches='tight', facecolor='white')

            return temp_file.name

        finally:
            plt.close(fig)  # CRITICAL: Prevent memory leaks

    def generate_horizon_distribution_chart(
        self,
        horizon_counts: Dict[str, int],
        title: str = "Horizon Distribution",
        dpi: int = CHART_DPI
    ) -> Optional[str]:
        """
        Generate a bar chart showing distribution of cards across horizons.

        Args:
            horizon_counts: Dictionary mapping horizon names to card counts
            title: Chart title
            dpi: Resolution for the image

        Returns:
            Path to the generated chart image, or None if no data
        """
        if not horizon_counts:
            logger.warning("No horizon data for distribution chart")
            return None

        fig, ax = plt.subplots(figsize=(6, 4))

        try:
            # Order horizons properly
            horizon_order = ['H1', 'H2', 'H3']
            labels = []
            values = []

            for h in horizon_order:
                if h in horizon_counts:
                    labels.append(h)
                    values.append(horizon_counts[h])

            # Add any remaining horizons
            for h, v in horizon_counts.items():
                if h not in horizon_order:
                    labels.append(h)
                    values.append(v)

            # Horizon colors
            horizon_colors = {
                'H1': FORESIGHT_COLORS["success"],
                'H2': FORESIGHT_COLORS["warning"],
                'H3': FORESIGHT_COLORS["secondary"],
            }
            colors = [horizon_colors.get(l, FORESIGHT_COLORS["primary"]) for l in labels]

            x_pos = np.arange(len(labels))
            bars = ax.bar(x_pos, values, color=colors, edgecolor='white', width=0.6)

            # Add value labels on bars
            for bar, value in zip(bars, values):
                height = bar.get_height()
                ax.text(
                    bar.get_x() + bar.get_width() / 2, height + 0.5,
                    f'{value}',
                    ha='center', va='bottom',
                    fontsize=11, fontweight='bold'
                )

            ax.set_xticks(x_pos)
            ax.set_xticklabels(labels, fontsize=12, fontweight='bold')
            ax.set_ylabel('Number of Cards', fontsize=11)
            ax.set_title(title, fontsize=12, fontweight='bold', pad=15)

            # Style
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.yaxis.grid(True, linestyle='--', alpha=0.3)
            ax.set_axisbelow(True)

            plt.tight_layout()

            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix='.png',
                delete=False,
                prefix='foresight_horizon_'
            )
            plt.savefig(temp_file.name, dpi=dpi, bbox_inches='tight', facecolor='white')

            return temp_file.name

        finally:
            plt.close(fig)  # CRITICAL: Prevent memory leaks

    # ========================================================================
    # PDF Generation Methods
    # ========================================================================

    def _get_pdf_styles(self) -> Dict[str, ParagraphStyle]:
        """
        Create custom paragraph styles for PDF generation.

        Returns:
            Dictionary of ParagraphStyle objects
        """
        styles = getSampleStyleSheet()

        custom_styles = {
            'Title': ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=PDF_TITLE_FONT_SIZE,
                textColor=PDF_COLORS["primary"],
                spaceAfter=12,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold',
            ),
            'Heading1': ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontSize=PDF_HEADING_FONT_SIZE,
                textColor=PDF_COLORS["primary"],
                spaceBefore=18,
                spaceAfter=8,
                fontName='Helvetica-Bold',
            ),
            'Heading2': ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=PDF_BODY_FONT_SIZE + 1,
                textColor=PDF_COLORS["secondary"],
                spaceBefore=12,
                spaceAfter=6,
                fontName='Helvetica-Bold',
            ),
            'Body': ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=PDF_BODY_FONT_SIZE,
                textColor=PDF_COLORS["dark"],
                spaceBefore=4,
                spaceAfter=4,
                leading=14,
            ),
            'Small': ParagraphStyle(
                'CustomSmall',
                parent=styles['Normal'],
                fontSize=PDF_SMALL_FONT_SIZE,
                textColor=rl_colors.gray,
                spaceBefore=2,
                spaceAfter=2,
            ),
            'Badge': ParagraphStyle(
                'Badge',
                parent=styles['Normal'],
                fontSize=PDF_SMALL_FONT_SIZE,
                textColor=rl_colors.white,
                alignment=TA_CENTER,
            ),
        }

        return custom_styles

    def _create_pdf_header(
        self,
        card_data: CardExportData,
        styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF header elements for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the header
        """
        elements = []

        # Title
        title = Paragraph(card_data.name, styles['Title'])
        elements.append(title)
        elements.append(Spacer(1, 6))

        # Horizontal rule
        elements.append(HRFlowable(
            width="100%",
            thickness=2,
            color=PDF_COLORS["primary"],
            spaceBefore=6,
            spaceAfter=12
        ))

        # Metadata badges (pillar, horizon, stage)
        badge_data = []
        if card_data.pillar_name or card_data.pillar_id:
            badge_data.append(('Pillar', card_data.pillar_name or card_data.pillar_id))
        if card_data.horizon:
            badge_data.append(('Horizon', card_data.horizon))
        if card_data.stage_name or card_data.stage_id:
            badge_data.append(('Stage', card_data.stage_name or card_data.stage_id))
        if card_data.goal_name or card_data.goal_id:
            badge_data.append(('Goal', card_data.goal_name or card_data.goal_id))

        if badge_data:
            badge_table_data = [[]]
            for label, value in badge_data:
                badge_text = f"<b>{label}:</b> {value}"
                badge_table_data[0].append(Paragraph(badge_text, styles['Small']))

            badge_table = Table(
                badge_table_data,
                colWidths=[1.5 * inch] * len(badge_data)
            )
            badge_table.setStyle(TableStyle([
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('TOPPADDING', (0, 0), (-1, -1), 4),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ]))
            elements.append(badge_table)
            elements.append(Spacer(1, 12))

        return elements

    def _create_pdf_summary_section(
        self,
        card_data: CardExportData,
        styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF summary section for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the summary
        """
        elements = []

        # Classification Overview - Always show this prominently
        elements.append(Paragraph("Classification Overview", styles['Heading1']))

        classification_items = []
        if card_data.pillar_name or card_data.pillar_id:
            classification_items.append(f"<b>Strategic Pillar:</b> {card_data.pillar_name or card_data.pillar_id}")
        if card_data.goal_name or card_data.goal_id:
            classification_items.append(f"<b>Strategic Goal:</b> {card_data.goal_name or card_data.goal_id}")
        if card_data.horizon:
            horizon_desc = {"H1": "Near-term (0-2 years)", "H2": "Mid-term (2-5 years)", "H3": "Long-term (5+ years)"}
            classification_items.append(f"<b>Time Horizon:</b> {card_data.horizon} - {horizon_desc.get(card_data.horizon, '')}")
        if card_data.stage_name or card_data.stage_id:
            classification_items.append(f"<b>Maturity Stage:</b> {card_data.stage_name or card_data.stage_id}")

        if classification_items:
            for item in classification_items:
                elements.append(Paragraph(item, styles['Body']))
            elements.append(Spacer(1, 12))

        # Executive Summary
        if card_data.summary:
            elements.append(Paragraph("Executive Summary", styles['Heading1']))
            elements.append(Paragraph(card_data.summary, styles['Body']))
            elements.append(Spacer(1, 12))

        # Full Description
        if card_data.description:
            elements.append(Paragraph("Detailed Analysis", styles['Heading1']))
            # Truncate very long descriptions
            description = card_data.description
            if len(description) > 3000:
                description = description[:3000] + "... [truncated]"
            elements.append(Paragraph(description, styles['Body']))
            elements.append(Spacer(1, 12))

        return elements

    def _create_pdf_research_section(
        self,
        card_data: CardExportData,
        styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF section for deep research report.

        Args:
            card_data: Card data containing research report
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the research section
        """
        elements = []

        if not card_data.deep_research_report:
            return elements

        elements.append(PageBreak())
        elements.append(Paragraph("Strategic Intelligence Report", styles['Title']))
        elements.append(Spacer(1, 6))

        # Add a note about the report
        elements.append(Paragraph(
            "<i>The following strategic intelligence report was generated through deep research analysis, "
            "synthesizing multiple sources to provide actionable insights for decision-makers.</i>",
            styles['Small']
        ))
        elements.append(Spacer(1, 12))

        elements.append(HRFlowable(
            width="100%",
            thickness=1,
            color=PDF_COLORS["secondary"],
            spaceBefore=6,
            spaceAfter=12
        ))

        # Parse and render the markdown report
        # Simple markdown-to-PDF conversion for common elements
        report = card_data.deep_research_report

        # Truncate if extremely long
        if len(report) > 15000:
            report = report[:15000] + "\n\n... [Report truncated for PDF export. View full report in the application.]"

        # Process the report line by line
        lines = report.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                elements.append(Spacer(1, 6))
            elif line.startswith('# '):
                # Main heading
                elements.append(Paragraph(line[2:], styles['Heading1']))
            elif line.startswith('## '):
                # Section heading
                elements.append(Paragraph(line[3:], styles['Heading2']))
            elif line.startswith('### '):
                # Subsection heading
                text = f"<b>{line[4:]}</b>"
                elements.append(Paragraph(text, styles['Body']))
            elif line.startswith('- ') or line.startswith('* '):
                # Bullet point
                text = f"• {line[2:]}"
                elements.append(Paragraph(text, styles['Body']))
            elif line.startswith('**') and line.endswith('**'):
                # Bold text
                text = f"<b>{line[2:-2]}</b>"
                elements.append(Paragraph(text, styles['Body']))
            else:
                # Regular paragraph
                # Escape any XML-unsafe characters
                safe_line = line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                elements.append(Paragraph(safe_line, styles['Body']))

        return elements

    def _create_pdf_scores_table(
        self,
        card_data: CardExportData,
        styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF scores table for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the scores section
        """
        elements = []

        elements.append(Paragraph("Scores", styles['Heading1']))

        # Build scores table
        scores = card_data.get_all_scores()
        table_data = [['Metric', 'Score', 'Rating']]

        for name, score in scores.items():
            score_display = self.format_score_display(score)

            # Determine rating based on score
            if score is None:
                rating = "Not Scored"
            elif score >= 80:
                rating = "Excellent"
            elif score >= 60:
                rating = "Good"
            elif score >= 40:
                rating = "Fair"
            else:
                rating = "Low"

            table_data.append([name, score_display, rating])

        # Create table with styling
        table = Table(table_data, colWidths=[2 * inch, 1 * inch, 1.5 * inch])
        table.setStyle(TableStyle([
            # Header styling
            ('BACKGROUND', (0, 0), (-1, 0), PDF_COLORS["primary"]),
            ('TEXTCOLOR', (0, 0), (-1, 0), rl_colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), PDF_BODY_FONT_SIZE),
            ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
            # Body styling
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), PDF_BODY_FONT_SIZE),
            ('ALIGN', (1, 1), (1, -1), 'CENTER'),  # Center score column
            ('ALIGN', (2, 1), (2, -1), 'CENTER'),  # Center rating column
            # Alternating row colors
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [rl_colors.white, PDF_COLORS["light"]]),
            # Grid
            ('GRID', (0, 0), (-1, -1), 0.5, PDF_COLORS["light"]),
            ('BOX', (0, 0), (-1, -1), 1, PDF_COLORS["primary"]),
            # Padding
            ('TOPPADDING', (0, 0), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('RIGHTPADDING', (0, 0), (-1, -1), 8),
        ]))

        elements.append(table)
        elements.append(Spacer(1, 18))

        return elements

    def _create_pdf_chart_section(
        self,
        card_data: CardExportData,
        styles: Dict[str, ParagraphStyle]
    ) -> Tuple[List[Any], List[str]]:
        """
        Create PDF chart section for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            Tuple of (list of flowable elements, list of temp file paths to clean up)
        """
        elements = []
        temp_files = []

        # Generate bar chart
        chart_path = self.generate_score_chart(card_data, chart_type="bar")

        if chart_path:
            temp_files.append(chart_path)
            elements.append(Paragraph("Score Visualization", styles['Heading1']))

            try:
                img = RLImage(chart_path, width=PDF_CHART_WIDTH, height=PDF_CHART_HEIGHT)
                elements.append(img)
                elements.append(Spacer(1, 12))
            except Exception as e:
                logger.warning(f"Failed to add chart image to PDF: {e}")

        return elements, temp_files

    def _create_pdf_footer(
        self,
        card_data: CardExportData,
        styles: Dict[str, ParagraphStyle]
    ) -> List[Any]:
        """
        Create PDF footer elements for a card.

        Args:
            card_data: Card data to display
            styles: PDF styles dictionary

        Returns:
            List of flowable elements for the footer
        """
        elements = []

        elements.append(Spacer(1, 24))
        elements.append(HRFlowable(
            width="100%",
            thickness=1,
            color=PDF_COLORS["light"],
            spaceBefore=6,
            spaceAfter=6
        ))

        # Metadata footer
        footer_parts = []
        if card_data.created_at:
            footer_parts.append(f"Created: {card_data.created_at.strftime('%Y-%m-%d')}")
        if card_data.updated_at:
            footer_parts.append(f"Updated: {card_data.updated_at.strftime('%Y-%m-%d')}")
        footer_parts.append(f"Export Date: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")

        footer_text = " | ".join(footer_parts)
        elements.append(Paragraph(footer_text, styles['Small']))

        # Branding
        elements.append(Spacer(1, 6))
        elements.append(Paragraph(
            "Generated by Foresight Intelligence Platform",
            styles['Small']
        ))

        return elements

    async def generate_pdf(
        self,
        card_data: CardExportData,
        include_charts: bool = True
    ) -> str:
        """
        Generate a PDF export for an intelligence card.

        Args:
            card_data: CardExportData object containing all card information
            include_charts: Whether to include chart visualizations

        Returns:
            Path to the generated PDF file

        Raises:
            Exception: If PDF generation fails
        """
        temp_files = []

        try:
            # Create temp file for PDF
            pdf_file = tempfile.NamedTemporaryFile(
                suffix='.pdf',
                delete=False,
                prefix='foresight_export_'
            )
            pdf_path = pdf_file.name
            pdf_file.close()

            # Create PDF document
            doc = SimpleDocTemplate(
                pdf_path,
                pagesize=PDF_PAGE_SIZE,
                rightMargin=PDF_MARGIN,
                leftMargin=PDF_MARGIN,
                topMargin=PDF_MARGIN,
                bottomMargin=PDF_MARGIN
            )

            # Get styles
            styles = self._get_pdf_styles()

            # Build document elements
            elements = []

            # Header with title and badges
            elements.extend(self._create_pdf_header(card_data, styles))

            # Summary and description
            elements.extend(self._create_pdf_summary_section(card_data, styles))

            # Scores table
            elements.extend(self._create_pdf_scores_table(card_data, styles))

            # Charts (if enabled)
            if include_charts:
                chart_elements, chart_files = self._create_pdf_chart_section(card_data, styles)
                elements.extend(chart_elements)
                temp_files.extend(chart_files)

            # Deep Research Report (if available)
            elements.extend(self._create_pdf_research_section(card_data, styles))

            # Footer
            elements.extend(self._create_pdf_footer(card_data, styles))

            # Build PDF
            doc.build(elements)

            logger.info(f"Generated PDF export for card: {card_data.name}")
            return pdf_path

        except Exception as e:
            logger.error(f"Error generating PDF for card {card_data.name}: {e}")
            raise

        finally:
            # Clean up chart temp files
            self.cleanup_temp_files(temp_files)

    async def generate_workstream_pdf(
        self,
        workstream_id: str,
        include_charts: bool = True,
        max_cards: int = 50
    ) -> str:
        """
        Generate a PDF report for a workstream containing all associated cards.

        Args:
            workstream_id: UUID of the workstream
            include_charts: Whether to include chart visualizations
            max_cards: Maximum number of cards to include

        Returns:
            Path to the generated PDF file

        Raises:
            Exception: If PDF generation fails
        """
        temp_files = []

        try:
            # Fetch workstream and cards
            workstream, cards = await self.get_workstream_cards(workstream_id, max_cards)

            if not workstream:
                raise ValueError(f"Workstream {workstream_id} not found")

            # Create temp file for PDF
            pdf_file = tempfile.NamedTemporaryFile(
                suffix='.pdf',
                delete=False,
                prefix='foresight_workstream_'
            )
            pdf_path = pdf_file.name
            pdf_file.close()

            # Create PDF document
            doc = SimpleDocTemplate(
                pdf_path,
                pagesize=PDF_PAGE_SIZE,
                rightMargin=PDF_MARGIN,
                leftMargin=PDF_MARGIN,
                topMargin=PDF_MARGIN,
                bottomMargin=PDF_MARGIN
            )

            # Get styles
            styles = self._get_pdf_styles()

            # Build document elements
            elements = []

            # Title page
            workstream_name = workstream.get('name', 'Workstream Report')
            elements.append(Paragraph(workstream_name, styles['Title']))
            elements.append(Spacer(1, 12))

            elements.append(HRFlowable(
                width="100%",
                thickness=2,
                color=PDF_COLORS["primary"],
                spaceBefore=6,
                spaceAfter=12
            ))

            # Workstream description
            if workstream.get('description'):
                elements.append(Paragraph("Overview", styles['Heading1']))
                elements.append(Paragraph(workstream['description'], styles['Body']))
                elements.append(Spacer(1, 12))

            # Summary statistics
            elements.append(Paragraph("Summary", styles['Heading1']))

            if not cards:
                elements.append(Paragraph(
                    "No cards currently match this workstream criteria.",
                    styles['Body']
                ))
            else:
                summary_text = f"This workstream contains <b>{len(cards)}</b> intelligence cards."
                if len(cards) >= max_cards:
                    summary_text += f" (Showing first {max_cards})"
                elements.append(Paragraph(summary_text, styles['Body']))
                elements.append(Spacer(1, 12))

                # Distribution charts
                if include_charts and cards:
                    # Pillar distribution
                    pillar_counts = {}
                    horizon_counts = {}
                    for card in cards:
                        pillar = card.pillar_name or card.pillar_id or "Unknown"
                        pillar_counts[pillar] = pillar_counts.get(pillar, 0) + 1

                        if card.horizon:
                            horizon_counts[card.horizon] = horizon_counts.get(card.horizon, 0) + 1

                    # Pillar chart
                    pillar_chart_path = self.generate_pillar_distribution_chart(pillar_counts)
                    if pillar_chart_path:
                        temp_files.append(pillar_chart_path)
                        try:
                            img = RLImage(pillar_chart_path, width=PDF_CHART_WIDTH, height=PDF_CHART_HEIGHT)
                            elements.append(img)
                            elements.append(Spacer(1, 12))
                        except Exception as e:
                            logger.warning(f"Failed to add pillar chart to PDF: {e}")

                    # Horizon chart
                    horizon_chart_path = self.generate_horizon_distribution_chart(horizon_counts)
                    if horizon_chart_path:
                        temp_files.append(horizon_chart_path)
                        try:
                            img = RLImage(horizon_chart_path, width=4.5 * inch, height=3 * inch)
                            elements.append(img)
                            elements.append(Spacer(1, 12))
                        except Exception as e:
                            logger.warning(f"Failed to add horizon chart to PDF: {e}")

                # Cards table summary
                elements.append(PageBreak())
                elements.append(Paragraph("Cards Overview", styles['Heading1']))

                table_data = [['Name', 'Pillar', 'Horizon', 'Impact']]
                for card in cards:
                    table_data.append([
                        card.name[:40] + ('...' if len(card.name) > 40 else ''),
                        card.pillar_name or card.pillar_id or 'N/A',
                        card.horizon or 'N/A',
                        self.format_score_display(card.impact_score)
                    ])

                table = Table(table_data, colWidths=[2.5 * inch, 1.5 * inch, 0.8 * inch, 0.8 * inch])
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), PDF_COLORS["primary"]),
                    ('TEXTCOLOR', (0, 0), (-1, 0), rl_colors.white),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), PDF_BODY_FONT_SIZE),
                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 1), (-1, -1), PDF_SMALL_FONT_SIZE),
                    ('ALIGN', (2, 0), (-1, -1), 'CENTER'),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [rl_colors.white, PDF_COLORS["light"]]),
                    ('GRID', (0, 0), (-1, -1), 0.5, PDF_COLORS["light"]),
                    ('BOX', (0, 0), (-1, -1), 1, PDF_COLORS["primary"]),
                    ('TOPPADDING', (0, 0), (-1, -1), 6),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ('LEFTPADDING', (0, 0), (-1, -1), 6),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                ]))
                elements.append(table)

                # Individual card details (paginated)
                elements.append(PageBreak())
                elements.append(Paragraph("Card Details", styles['Heading1']))
                elements.append(Spacer(1, 12))

                for i, card in enumerate(cards):
                    if i > 0:
                        elements.append(Spacer(1, 18))
                        elements.append(HRFlowable(
                            width="80%",
                            thickness=0.5,
                            color=PDF_COLORS["light"],
                            spaceBefore=6,
                            spaceAfter=12
                        ))

                    # Card mini-header
                    elements.append(Paragraph(card.name, styles['Heading2']))

                    # Badges row
                    badge_parts = []
                    if card.pillar_name or card.pillar_id:
                        badge_parts.append(f"<b>Pillar:</b> {card.pillar_name or card.pillar_id}")
                    if card.horizon:
                        badge_parts.append(f"<b>Horizon:</b> {card.horizon}")
                    if card.stage_name or card.stage_id:
                        badge_parts.append(f"<b>Stage:</b> {card.stage_name or card.stage_id}")
                    if badge_parts:
                        elements.append(Paragraph(" | ".join(badge_parts), styles['Small']))
                        elements.append(Spacer(1, 6))

                    # Summary
                    if card.summary:
                        elements.append(Paragraph(card.summary, styles['Body']))

                    # Key scores
                    key_scores = []
                    if card.impact_score is not None:
                        key_scores.append(f"Impact: {card.impact_score}")
                    if card.relevance_score is not None:
                        key_scores.append(f"Relevance: {card.relevance_score}")
                    if card.maturity_score is not None:
                        key_scores.append(f"Maturity: {card.maturity_score}")
                    if key_scores:
                        elements.append(Paragraph(
                            "<i>Scores: " + " | ".join(key_scores) + "</i>",
                            styles['Small']
                        ))

            # Footer
            elements.append(Spacer(1, 24))
            elements.append(HRFlowable(
                width="100%",
                thickness=1,
                color=PDF_COLORS["light"],
                spaceBefore=6,
                spaceAfter=6
            ))

            footer_text = f"Export Date: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
            elements.append(Paragraph(footer_text, styles['Small']))
            elements.append(Paragraph(
                "Generated by Foresight Intelligence Platform",
                styles['Small']
            ))

            # Build PDF
            doc.build(elements)

            logger.info(f"Generated workstream PDF for: {workstream_name} with {len(cards)} cards")
            return pdf_path

        except Exception as e:
            logger.error(f"Error generating workstream PDF {workstream_id}: {e}")
            raise

        finally:
            # Clean up chart temp files
            self.cleanup_temp_files(temp_files)

    # ========================================================================
    # Utility Methods
    # ========================================================================

    def cleanup_temp_files(self, file_paths: List[str]) -> None:
        """
        Clean up temporary chart files.

        Args:
            file_paths: List of file paths to delete
        """
        for path in file_paths:
            try:
                if path and Path(path).exists():
                    Path(path).unlink()
                    logger.debug(f"Cleaned up temp file: {path}")
            except Exception as e:
                logger.warning(f"Failed to clean up temp file {path}: {e}")

    async def get_card_data(self, card_id: str) -> Optional[CardExportData]:
        """
        Fetch card data from database and convert to CardExportData.

        Args:
            card_id: UUID of the card to fetch

        Returns:
            CardExportData object or None if not found
        """
        try:
            response = self.supabase.table("cards").select("*").eq("id", card_id).single().execute()

            if not response.data:
                return None

            return CardExportData(**response.data)

        except Exception as e:
            logger.error(f"Error fetching card {card_id}: {e}")
            return None

    async def get_workstream_cards(
        self,
        workstream_id: str,
        max_cards: int = 50
    ) -> Tuple[Optional[Dict[str, Any]], List[CardExportData]]:
        """
        Fetch workstream metadata and associated cards.

        Args:
            workstream_id: UUID of the workstream
            max_cards: Maximum number of cards to fetch

        Returns:
            Tuple of (workstream_data, list of CardExportData)
        """
        try:
            # Fetch workstream
            ws_response = self.supabase.table("workstreams").select("*").eq("id", workstream_id).single().execute()

            if not ws_response.data:
                return None, []

            workstream = ws_response.data

            # Fetch associated cards via workstream_cards junction table
            cards_response = self.supabase.table("workstream_cards").select(
                "card_id, cards(*)"
            ).eq("workstream_id", workstream_id).limit(max_cards).execute()

            cards = []
            if cards_response.data:
                for item in cards_response.data:
                    if item.get("cards"):
                        cards.append(CardExportData(**item["cards"]))

            return workstream, cards

        except Exception as e:
            logger.error(f"Error fetching workstream {workstream_id}: {e}")
            return None, []

    def format_score_display(self, score: Optional[int]) -> str:
        """
        Format a score for display, handling None values.

        Args:
            score: Score value (0-100) or None

        Returns:
            Formatted string representation
        """
        return str(score) if score is not None else "N/A"

    def get_content_type(self, format: ExportFormat) -> str:
        """
        Get the MIME content type for an export format.

        Args:
            format: Export format

        Returns:
            MIME content type string
        """
        return EXPORT_CONTENT_TYPES.get(format, "application/octet-stream")

    def generate_filename(self, name: str, format: ExportFormat) -> str:
        """
        Generate a safe filename for an export.

        Args:
            name: Card or workstream name
            format: Export format

        Returns:
            Safe filename with extension
        """
        return get_export_filename(name, format)

    # ========================================================================
    # CSV Export Methods
    # ========================================================================

    async def generate_csv(
        self,
        card_data: CardExportData,
    ) -> str:
        """
        Generate CSV export for a single intelligence card.

        Exports card data in a tabular format suitable for analysis
        in Excel or other spreadsheet applications. All card fields
        and scores are included as columns.

        Args:
            card_data: Card data to export

        Returns:
            CSV string content (not file path)

        Raises:
            ValueError: If card_data is invalid
        """
        import pandas as pd

        try:
            # Define the CSV columns in the specified order
            csv_columns = [
                "id",
                "name",
                "summary",
                "description",
                "pillar_id",
                "goal_id",
                "stage_id",
                "horizon",
                "novelty_score",
                "maturity_score",
                "impact_score",
                "relevance_score",
                "velocity_score",
                "risk_score",
                "opportunity_score",
            ]

            # Build the row data from card_data
            row_data = {
                "id": card_data.id,
                "name": card_data.name,
                "summary": card_data.summary or "",
                "description": card_data.description or "",
                "pillar_id": card_data.pillar_id or "",
                "goal_id": card_data.goal_id or "",
                "stage_id": card_data.stage_id or "",
                "horizon": card_data.horizon or "",
                "novelty_score": card_data.novelty_score,
                "maturity_score": card_data.maturity_score,
                "impact_score": card_data.impact_score,
                "relevance_score": card_data.relevance_score,
                "velocity_score": card_data.velocity_score,
                "risk_score": card_data.risk_score,
                "opportunity_score": card_data.opportunity_score,
            }

            # Create DataFrame with single row
            df = pd.DataFrame([row_data], columns=csv_columns)

            # Convert to CSV string without index column
            csv_content = df.to_csv(index=False)

            logger.info(f"Generated CSV export for card {card_data.id}: {card_data.name}")

            return csv_content

        except Exception as e:
            logger.error(f"Error generating CSV for card {card_data.id}: {e}")
            raise ValueError(f"Failed to generate CSV export: {e}")

    async def generate_csv_multi(
        self,
        cards: List[CardExportData],
    ) -> str:
        """
        Generate CSV export for multiple intelligence cards.

        Exports multiple cards as rows in a single CSV file,
        suitable for bulk data analysis in Excel or other tools.

        Args:
            cards: List of card data to export

        Returns:
            CSV string content with multiple rows

        Raises:
            ValueError: If cards list is empty or invalid
        """
        import pandas as pd

        if not cards:
            logger.warning("No cards provided for CSV export")
            return self._generate_empty_csv()

        try:
            # Define the CSV columns in the specified order
            csv_columns = [
                "id",
                "name",
                "summary",
                "description",
                "pillar_id",
                "goal_id",
                "stage_id",
                "horizon",
                "novelty_score",
                "maturity_score",
                "impact_score",
                "relevance_score",
                "velocity_score",
                "risk_score",
                "opportunity_score",
            ]

            # Build row data for all cards
            rows = []
            for card_data in cards:
                row = {
                    "id": card_data.id,
                    "name": card_data.name,
                    "summary": card_data.summary or "",
                    "description": card_data.description or "",
                    "pillar_id": card_data.pillar_id or "",
                    "goal_id": card_data.goal_id or "",
                    "stage_id": card_data.stage_id or "",
                    "horizon": card_data.horizon or "",
                    "novelty_score": card_data.novelty_score,
                    "maturity_score": card_data.maturity_score,
                    "impact_score": card_data.impact_score,
                    "relevance_score": card_data.relevance_score,
                    "velocity_score": card_data.velocity_score,
                    "risk_score": card_data.risk_score,
                    "opportunity_score": card_data.opportunity_score,
                }
                rows.append(row)

            # Create DataFrame with all rows
            df = pd.DataFrame(rows, columns=csv_columns)

            # Convert to CSV string without index column
            csv_content = df.to_csv(index=False)

            logger.info(f"Generated CSV export for {len(cards)} cards")

            return csv_content

        except Exception as e:
            logger.error(f"Error generating multi-card CSV: {e}")
            raise ValueError(f"Failed to generate CSV export: {e}")

    def _generate_empty_csv(self) -> str:
        """
        Generate an empty CSV with just headers.

        Returns:
            CSV string with headers only
        """
        import pandas as pd

        csv_columns = [
            "id",
            "name",
            "summary",
            "description",
            "pillar_id",
            "goal_id",
            "stage_id",
            "horizon",
            "novelty_score",
            "maturity_score",
            "impact_score",
            "relevance_score",
            "velocity_score",
            "risk_score",
            "opportunity_score",
        ]

        df = pd.DataFrame(columns=csv_columns)
        return df.to_csv(index=False)

    # ========================================================================
    # PowerPoint Export Methods
    # ========================================================================

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """
        Convert hex color string to RGBColor for PowerPoint.

        Args:
            hex_color: Hex color string (e.g., '#1E3A5F')

        Returns:
            RGBColor object for use with python-pptx
        """
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def _add_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: Optional[str] = None
    ) -> None:
        """
        Add a title slide to the presentation.

        Args:
            prs: Presentation object
            title: Main title text
            subtitle: Optional subtitle text
        """
        # Use blank layout for custom styling
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Add background color shape
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            PPTX_SLIDE_WIDTH, PPTX_SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["primary"])
        background.line.fill.background()

        # Add title text box
        title_box = slide.shapes.add_textbox(
            PPTX_MARGIN, Inches(2.5),
            PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title[:80]  # Truncate long titles
        title_para.font.size = PPTX_TITLE_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                PPTX_MARGIN, Inches(4.2),
                PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle[:150]
            subtitle_para.font.size = PPTX_SUBTITLE_FONT_SIZE
            subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Add Foresight branding footer
        footer_box = slide.shapes.add_textbox(
            PPTX_MARGIN, PPTX_SLIDE_HEIGHT - Inches(0.8),
            PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.4)
        )
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.paragraphs[0]
        footer_para.text = f"Foresight Intelligence Platform • {datetime.now().strftime('%B %d, %Y')}"
        footer_para.font.size = PPTX_SMALL_FONT_SIZE
        footer_para.font.color.rgb = RGBColor(180, 180, 180)
        footer_para.alignment = PP_ALIGN.CENTER

    def _add_content_slide(
        self,
        prs: Presentation,
        title: str,
        content_items: List[Tuple[str, str]],
        chart_path: Optional[str] = None
    ) -> None:
        """
        Add a content slide with text and optional chart.

        Args:
            prs: Presentation object
            title: Slide title
            content_items: List of (label, value) tuples
            chart_path: Optional path to chart image to include
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            PPTX_SLIDE_WIDTH, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["primary"])
        title_bar.line.fill.background()

        # Add title text
        title_box = slide.shapes.add_textbox(
            PPTX_MARGIN, Inches(0.3),
            PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title[:60]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Determine layout based on whether chart is included
        if chart_path:
            content_width = Inches(6.5)
            chart_left = Inches(7.5)
        else:
            content_width = PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN)
            chart_left = None

        # Add content items
        content_top = Inches(1.6)
        content_box = slide.shapes.add_textbox(
            PPTX_MARGIN, content_top,
            content_width, Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, (label, value) in enumerate(content_items):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()

            para.space_before = Pt(8)
            para.space_after = Pt(4)

            # Add label in bold
            run_label = para.add_run()
            run_label.text = f"{label}: "
            run_label.font.size = PPTX_BODY_FONT_SIZE
            run_label.font.bold = True
            run_label.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])

            # Add value
            run_value = para.add_run()
            run_value.text = str(value) if value else "N/A"
            run_value.font.size = PPTX_BODY_FONT_SIZE
            run_value.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])

        # Add chart if provided
        if chart_path and Path(chart_path).exists():
            try:
                slide.shapes.add_picture(
                    chart_path,
                    chart_left, Inches(1.8),
                    width=PPTX_CHART_WIDTH, height=PPTX_CHART_HEIGHT
                )
            except Exception as e:
                logger.warning(f"Failed to add chart to slide: {e}")

    def _add_scores_slide(
        self,
        prs: Presentation,
        card_data: CardExportData,
        chart_path: Optional[str] = None
    ) -> None:
        """
        Add a slide showing all scores with optional chart.

        Args:
            prs: Presentation object
            card_data: Card data with scores
            chart_path: Optional path to score chart image
        """
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            PPTX_SLIDE_WIDTH, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["secondary"])
        title_bar.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            PPTX_MARGIN, Inches(0.3),
            PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Score Analysis"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add chart if available
        if chart_path and Path(chart_path).exists():
            try:
                slide.shapes.add_picture(
                    chart_path,
                    Inches(0.5), Inches(1.5),
                    width=Inches(6), height=Inches(5)
                )
            except Exception as e:
                logger.warning(f"Failed to add score chart: {e}")

        # Add score details on the right side
        scores = card_data.get_all_scores()
        scores_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.8),
            Inches(5.5), Inches(5)
        )
        scores_frame = scores_box.text_frame
        scores_frame.word_wrap = True

        for i, (score_name, score_value) in enumerate(scores.items()):
            if i == 0:
                para = scores_frame.paragraphs[0]
            else:
                para = scores_frame.add_paragraph()

            para.space_before = Pt(12)
            para.space_after = Pt(4)

            # Score name
            run_name = para.add_run()
            run_name.text = f"{score_name}: "
            run_name.font.size = Pt(20)
            run_name.font.bold = True
            run_name.font.color.rgb = self._hex_to_rgb(
                SCORE_COLORS.get(score_name, FORESIGHT_COLORS["dark"])
            )

            # Score value
            run_value = para.add_run()
            run_value.text = str(score_value) if score_value is not None else "N/A"
            run_value.font.size = Pt(20)
            run_value.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])

    def _add_description_slide(
        self,
        prs: Presentation,
        title: str,
        description: Optional[str]
    ) -> None:
        """
        Add a slide for long-form description text.

        Args:
            prs: Presentation object
            title: Slide title
            description: Description text (will be truncated if too long)
        """
        if not description:
            return

        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            PPTX_SLIDE_WIDTH, Inches(1.2)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["primary"])
        title_bar.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            PPTX_MARGIN, Inches(0.3),
            PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add description text - truncate if too long for slide
        max_chars = 2000  # Reasonable limit for one slide
        display_text = description[:max_chars]
        if len(description) > max_chars:
            display_text += "..."

        desc_box = slide.shapes.add_textbox(
            PPTX_MARGIN, Inches(1.6),
            PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(5.5)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = display_text
        desc_para.font.size = PPTX_BODY_FONT_SIZE
        desc_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])
        desc_para.line_spacing = 1.3

    async def generate_pptx(
        self,
        card_data: CardExportData,
        include_charts: bool = True,
        include_description: bool = True
    ) -> str:
        """
        Generate a PowerPoint presentation for an individual card.

        Creates a multi-slide presentation with:
        - Title slide with card name and summary
        - Overview slide with key metadata
        - Scores slide with visualization
        - Description slide (if enabled)

        Args:
            card_data: Card data to export
            include_charts: Whether to include score charts
            include_description: Whether to include description slide

        Returns:
            Path to the generated PowerPoint file

        Raises:
            Exception: If PowerPoint generation fails
        """
        temp_files_to_cleanup = []

        try:
            logger.info(f"Generating PowerPoint for card: {card_data.name}")

            # Create presentation
            prs = Presentation()
            prs.slide_width = PPTX_SLIDE_WIDTH
            prs.slide_height = PPTX_SLIDE_HEIGHT

            # 1. Title slide
            self._add_title_slide(
                prs,
                title=card_data.name,
                subtitle=card_data.summary
            )

            # 2. Overview slide with metadata
            overview_items = [
                ("Pillar", card_data.pillar_name or card_data.pillar_id),
                ("Goal", card_data.goal_name or card_data.goal_id),
                ("Anchor", card_data.anchor_name or card_data.anchor_id),
                ("Stage", card_data.stage_name or card_data.stage_id),
                ("Horizon", card_data.horizon),
                ("Status", card_data.status),
            ]
            # Filter out items with no value
            overview_items = [(k, v) for k, v in overview_items if v]

            self._add_content_slide(
                prs,
                title="Card Overview",
                content_items=overview_items
            )

            # 3. Scores slide with chart
            chart_path = None
            if include_charts:
                chart_path = self.generate_score_chart(card_data, chart_type="radar")
                if chart_path:
                    temp_files_to_cleanup.append(chart_path)

            self._add_scores_slide(prs, card_data, chart_path)

            # 4. Description slide (optional)
            if include_description and card_data.description:
                self._add_description_slide(
                    prs,
                    title="Full Description",
                    description=card_data.description
                )

            # Save presentation to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix='.pptx',
                delete=False,
                prefix='foresight_card_'
            )
            prs.save(temp_file.name)

            logger.info(f"PowerPoint generated successfully: {temp_file.name}")
            return temp_file.name

        except Exception as e:
            logger.error(f"Error generating PowerPoint: {e}")
            raise

        finally:
            # Clean up chart temp files
            self.cleanup_temp_files(temp_files_to_cleanup)

    async def generate_workstream_pptx(
        self,
        workstream: Dict[str, Any],
        cards: List[CardExportData],
        include_charts: bool = True,
        include_card_details: bool = True
    ) -> str:
        """
        Generate a PowerPoint presentation for a workstream report.

        Creates a comprehensive presentation with:
        - Title slide with workstream name
        - Summary slide with statistics
        - Distribution charts (pillar, horizon)
        - Individual card slides (if enabled)

        Args:
            workstream: Workstream metadata dict
            cards: List of cards in the workstream
            include_charts: Whether to include distribution charts
            include_card_details: Whether to include individual card slides

        Returns:
            Path to the generated PowerPoint file

        Raises:
            Exception: If PowerPoint generation fails
        """
        temp_files_to_cleanup = []

        try:
            workstream_name = workstream.get('name', 'Workstream Report')
            logger.info(f"Generating workstream PowerPoint: {workstream_name}")

            # Create presentation
            prs = Presentation()
            prs.slide_width = PPTX_SLIDE_WIDTH
            prs.slide_height = PPTX_SLIDE_HEIGHT

            # 1. Title slide
            self._add_title_slide(
                prs,
                title=workstream_name,
                subtitle=f"Intelligence Report • {len(cards)} Cards"
            )

            # 2. Summary slide
            summary_items = [
                ("Total Cards", str(len(cards))),
                ("Description", workstream.get('description', 'N/A')),
            ]

            # Calculate pillar distribution
            pillar_counts: Dict[str, int] = {}
            horizon_counts: Dict[str, int] = {}
            for card in cards:
                pillar = card.pillar_name or card.pillar_id or "Unknown"
                pillar_counts[pillar] = pillar_counts.get(pillar, 0) + 1

                horizon = card.horizon or "Unknown"
                horizon_counts[horizon] = horizon_counts.get(horizon, 0) + 1

            if pillar_counts:
                pillar_summary = ", ".join(f"{k}: {v}" for k, v in pillar_counts.items())
                summary_items.append(("Pillars", pillar_summary))

            if horizon_counts:
                horizon_summary = ", ".join(f"{k}: {v}" for k, v in horizon_counts.items())
                summary_items.append(("Horizons", horizon_summary))

            self._add_content_slide(
                prs,
                title="Workstream Summary",
                content_items=summary_items
            )

            # 3. Distribution charts slide
            if include_charts and cards:
                # Generate pillar distribution chart
                pillar_chart_path = None
                if pillar_counts:
                    pillar_chart_path = self.generate_pillar_distribution_chart(pillar_counts)
                    if pillar_chart_path:
                        temp_files_to_cleanup.append(pillar_chart_path)

                # Generate horizon distribution chart
                horizon_chart_path = None
                if horizon_counts:
                    horizon_chart_path = self.generate_horizon_distribution_chart(horizon_counts)
                    if horizon_chart_path:
                        temp_files_to_cleanup.append(horizon_chart_path)

                # Add distribution slide with both charts
                if pillar_chart_path or horizon_chart_path:
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)

                    # Title bar
                    title_bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0),
                        PPTX_SLIDE_WIDTH, Inches(1.2)
                    )
                    title_bar.fill.solid()
                    title_bar.fill.fore_color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["accent"])
                    title_bar.line.fill.background()

                    title_box = slide.shapes.add_textbox(
                        PPTX_MARGIN, Inches(0.3),
                        PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.8)
                    )
                    title_frame = title_box.text_frame
                    title_para = title_frame.paragraphs[0]
                    title_para.text = "Distribution Analysis"
                    title_para.font.size = Pt(32)
                    title_para.font.bold = True
                    title_para.font.color.rgb = RGBColor(255, 255, 255)

                    # Add pillar chart on left
                    if pillar_chart_path and Path(pillar_chart_path).exists():
                        try:
                            slide.shapes.add_picture(
                                pillar_chart_path,
                                Inches(0.3), Inches(1.5),
                                width=Inches(6), height=Inches(5)
                            )
                        except Exception as e:
                            logger.warning(f"Failed to add pillar chart: {e}")

                    # Add horizon chart on right
                    if horizon_chart_path and Path(horizon_chart_path).exists():
                        try:
                            slide.shapes.add_picture(
                                horizon_chart_path,
                                Inches(7), Inches(1.5),
                                width=Inches(5.5), height=Inches(4.5)
                            )
                        except Exception as e:
                            logger.warning(f"Failed to add horizon chart: {e}")

            # 4. Individual card slides
            if include_card_details:
                # Handle empty workstream case
                if not cards:
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)

                    msg_box = slide.shapes.add_textbox(
                        Inches(2), Inches(3),
                        Inches(9), Inches(2)
                    )
                    msg_frame = msg_box.text_frame
                    msg_para = msg_frame.paragraphs[0]
                    msg_para.text = "No cards currently match this workstream criteria"
                    msg_para.font.size = PPTX_SUBTITLE_FONT_SIZE
                    msg_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])
                    msg_para.alignment = PP_ALIGN.CENTER
                else:
                    # Add a slide for each card (up to 50)
                    for card in cards[:50]:
                        card_items = [
                            ("Summary", card.summary),
                            ("Pillar", card.pillar_name or card.pillar_id),
                            ("Horizon", card.horizon),
                            ("Stage", card.stage_name or card.stage_id),
                        ]
                        # Add scores
                        scores = card.get_all_scores()
                        valid_scores = {k: v for k, v in scores.items() if v is not None}
                        if valid_scores:
                            scores_text = ", ".join(f"{k}: {v}" for k, v in valid_scores.items())
                            card_items.append(("Scores", scores_text))

                        # Filter out empty items
                        card_items = [(k, v) for k, v in card_items if v]

                        self._add_content_slide(
                            prs,
                            title=card.name[:50],
                            content_items=card_items
                        )

            # Save presentation to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix='.pptx',
                delete=False,
                prefix='foresight_workstream_'
            )
            prs.save(temp_file.name)

            logger.info(f"Workstream PowerPoint generated: {temp_file.name}")
            return temp_file.name

        except Exception as e:
            logger.error(f"Error generating workstream PowerPoint: {e}")
            raise

        finally:
            # Clean up chart temp files
            self.cleanup_temp_files(temp_files_to_cleanup)

    # ========================================================================
    # Executive Brief Export Methods
    # ========================================================================

    async def generate_brief_pdf(
        self,
        brief_title: str,
        card_name: str,
        executive_summary: str,
        content_markdown: str,
        generated_at: Optional[datetime] = None,
        version: Optional[int] = None
    ) -> str:
        """
        Generate a PDF export for an executive brief.

        Args:
            brief_title: Title for the brief (usually card name)
            card_name: Name of the card this brief is for
            executive_summary: Executive summary text
            content_markdown: Full brief content in markdown format
            generated_at: When the brief was generated
            version: Version number of the brief

        Returns:
            Path to the generated PDF file

        Raises:
            Exception: If PDF generation fails
        """
        try:
            # Create temp file for PDF
            pdf_file = tempfile.NamedTemporaryFile(
                suffix='.pdf',
                delete=False,
                prefix='foresight_brief_'
            )
            pdf_path = pdf_file.name
            pdf_file.close()

            # Create PDF document
            doc = SimpleDocTemplate(
                pdf_path,
                pagesize=PDF_PAGE_SIZE,
                rightMargin=PDF_MARGIN,
                leftMargin=PDF_MARGIN,
                topMargin=PDF_MARGIN,
                bottomMargin=PDF_MARGIN
            )

            # Get styles
            styles = self._get_pdf_styles()

            # Build document elements
            elements = []

            # Title
            title_text = f"Executive Brief: {brief_title}"
            if version and version > 1:
                title_text += f" (v{version})"
            elements.append(Paragraph(title_text, styles['Title']))
            elements.append(Spacer(1, 6))

            # Horizontal rule
            elements.append(HRFlowable(
                width="100%",
                thickness=2,
                color=PDF_COLORS["primary"],
                spaceBefore=6,
                spaceAfter=12
            ))

            # Metadata
            meta_parts = [f"Card: {card_name}"]
            if generated_at:
                meta_parts.append(f"Generated: {generated_at.strftime('%Y-%m-%d %H:%M UTC')}")
            if version:
                meta_parts.append(f"Version: {version}")
            elements.append(Paragraph(" | ".join(meta_parts), styles['Small']))
            elements.append(Spacer(1, 12))

            # Executive Summary section
            elements.append(Paragraph("Executive Summary", styles['Heading1']))
            elements.append(Paragraph(executive_summary or "No summary available.", styles['Body']))
            elements.append(Spacer(1, 18))

            # Full Brief Content
            elements.append(Paragraph("Full Brief", styles['Heading1']))
            elements.append(Spacer(1, 6))

            # Parse markdown content into paragraphs
            # Simple markdown parsing - split by double newlines for paragraphs
            if content_markdown:
                paragraphs = content_markdown.split('\n\n')
                for para_text in paragraphs:
                    para_text = para_text.strip()
                    if not para_text:
                        continue

                    # Handle headers
                    if para_text.startswith('# '):
                        elements.append(Paragraph(para_text[2:], styles['Heading1']))
                    elif para_text.startswith('## '):
                        elements.append(Paragraph(para_text[3:], styles['Heading2']))
                    elif para_text.startswith('### '):
                        # Create a heading3 style on the fly
                        heading3_style = ParagraphStyle(
                            'Heading3',
                            parent=styles['Body'],
                            fontSize=PDF_BODY_FONT_SIZE,
                            fontName='Helvetica-Bold',
                            spaceBefore=10,
                            spaceAfter=4,
                        )
                        elements.append(Paragraph(para_text[4:], heading3_style))
                    elif para_text.startswith('- ') or para_text.startswith('* '):
                        # Bullet list items
                        bullet_style = ParagraphStyle(
                            'Bullet',
                            parent=styles['Body'],
                            leftIndent=20,
                            firstLineIndent=-10,
                            spaceBefore=2,
                            spaceAfter=2,
                        )
                        # Handle multi-line bullets
                        lines = para_text.split('\n')
                        for line in lines:
                            line = line.strip()
                            if line.startswith('- ') or line.startswith('* '):
                                elements.append(Paragraph(f"• {line[2:]}", bullet_style))
                            elif line:
                                elements.append(Paragraph(line, styles['Body']))
                    else:
                        # Regular paragraph - handle inline markdown
                        # Convert **bold** to <b>bold</b>
                        import re
                        formatted = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', para_text)
                        # Convert *italic* to <i>italic</i>
                        formatted = re.sub(r'\*(.+?)\*', r'<i>\1</i>', formatted)
                        # Handle line breaks within paragraph
                        formatted = formatted.replace('\n', '<br/>')
                        elements.append(Paragraph(formatted, styles['Body']))

                    elements.append(Spacer(1, 4))
            else:
                elements.append(Paragraph("No content available.", styles['Body']))

            # Footer
            elements.append(Spacer(1, 24))
            elements.append(HRFlowable(
                width="100%",
                thickness=1,
                color=PDF_COLORS["light"],
                spaceBefore=6,
                spaceAfter=6
            ))

            footer_text = f"Export Date: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
            elements.append(Paragraph(footer_text, styles['Small']))
            elements.append(Paragraph(
                "Generated by Foresight Intelligence Platform",
                styles['Small']
            ))

            # Build PDF
            doc.build(elements)

            logger.info(f"Generated brief PDF export: {brief_title}")
            return pdf_path

        except Exception as e:
            logger.error(f"Error generating brief PDF: {e}")
            raise

    async def generate_professional_brief_pdf(
        self,
        brief_title: str,
        card_name: str,
        executive_summary: str,
        content_markdown: str,
        generated_at: Optional[datetime] = None,
        version: Optional[int] = None,
        classification: Optional[Dict[str, str]] = None
    ) -> str:
        """
        Generate a professional PDF export for an executive brief.
        
        This version includes:
        - City of Austin logo in header
        - Foresight Strategic Research branding
        - Professional header/footer on every page
        - Full AI technology disclosure
        
        Designed for senior city leadership distribution.

        Args:
            brief_title: Title for the brief
            card_name: Name of the card this brief is for
            executive_summary: Executive summary text
            content_markdown: Full brief content in markdown format
            generated_at: When the brief was generated
            version: Version number of the brief
            classification: Optional dict with pillar, horizon, stage info

        Returns:
            Path to the generated PDF file

        Raises:
            Exception: If PDF generation fails
        """
        import re
        
        try:
            # Create temp file for PDF
            pdf_file = tempfile.NamedTemporaryFile(
                suffix='.pdf',
                delete=False,
                prefix='foresight_executive_brief_'
            )
            pdf_path = pdf_file.name
            pdf_file.close()

            # Get professional styles
            styles = get_professional_pdf_styles()

            # Build document elements
            elements = []

            # Document Title
            elements.append(Paragraph(brief_title, styles['DocTitle']))
            
            # Subtitle with metadata
            subtitle_parts = []
            if classification:
                if classification.get('pillar'):
                    subtitle_parts.append(f"Pillar: {classification['pillar']}")
                if classification.get('horizon'):
                    subtitle_parts.append(f"Horizon: {classification['horizon']}")
                if classification.get('stage'):
                    subtitle_parts.append(f"Stage: {classification['stage']}")
            if subtitle_parts:
                elements.append(Paragraph(" | ".join(subtitle_parts), styles['DocSubtitle']))
            
            elements.append(Spacer(1, 6))
            
            # Decorative line
            elements.append(HRFlowable(
                width="100%",
                thickness=2,
                color=PDF_COLORS["secondary"],
                spaceBefore=4,
                spaceAfter=16
            ))

            # Executive Summary Section (highlighted)
            elements.append(Paragraph("Executive Summary", styles['SectionHeading']))
            
            if executive_summary:
                # Create a highlighted box for executive summary
                summary_text = executive_summary.replace('\n', '<br/>')
                elements.append(Paragraph(summary_text, styles['ExecutiveSummary']))
            else:
                elements.append(Paragraph("No summary available.", styles['BodyText']))
            
            elements.append(Spacer(1, 16))

            # Main Content Section
            elements.append(Paragraph("Strategic Intelligence Report", styles['SectionHeading']))
            elements.append(Spacer(1, 8))

            # Parse and render markdown content
            if content_markdown:
                lines = content_markdown.split('\n')
                current_paragraph = []
                
                for line in lines:
                    line_stripped = line.strip()
                    
                    # Empty line - flush current paragraph
                    if not line_stripped:
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            # Convert markdown formatting
                            para_text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', para_text)
                            para_text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', para_text)
                            # Escape XML characters
                            para_text = para_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            # Restore our formatting tags
                            para_text = para_text.replace('&lt;b&gt;', '<b>').replace('&lt;/b&gt;', '</b>')
                            para_text = para_text.replace('&lt;i&gt;', '<i>').replace('&lt;/i&gt;', '</i>')
                            elements.append(Paragraph(para_text, styles['BodyText']))
                            current_paragraph = []
                        continue
                    
                    # Headers
                    if line_stripped.startswith('# '):
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            elements.append(Paragraph(para_text, styles['BodyText']))
                            current_paragraph = []
                        elements.append(Spacer(1, 8))
                        elements.append(Paragraph(line_stripped[2:], styles['SectionHeading']))
                    elif line_stripped.startswith('## '):
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            elements.append(Paragraph(para_text, styles['BodyText']))
                            current_paragraph = []
                        elements.append(Paragraph(line_stripped[3:], styles['SubsectionHeading']))
                    elif line_stripped.startswith('### '):
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            elements.append(Paragraph(para_text, styles['BodyText']))
                            current_paragraph = []
                        elements.append(Paragraph(f"<b>{line_stripped[4:]}</b>", styles['BodyText']))
                    # Bullet points
                    elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            elements.append(Paragraph(para_text, styles['BodyText']))
                            current_paragraph = []
                        bullet_text = line_stripped[2:]
                        # Convert markdown in bullet
                        bullet_text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', bullet_text)
                        elements.append(Paragraph(f"• {bullet_text}", styles['BulletText']))
                    # Numbered lists
                    elif re.match(r'^\d+\.\s', line_stripped):
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            elements.append(Paragraph(para_text, styles['BodyText']))
                            current_paragraph = []
                        elements.append(Paragraph(line_stripped, styles['BulletText']))
                    # Horizontal rule
                    elif line_stripped == '---' or line_stripped == '***':
                        if current_paragraph:
                            para_text = ' '.join(current_paragraph)
                            elements.append(Paragraph(para_text, styles['BodyText']))
                            current_paragraph = []
                        elements.append(Spacer(1, 8))
                        elements.append(HRFlowable(
                            width="60%",
                            thickness=1,
                            color=PDF_COLORS["light"],
                            spaceBefore=4,
                            spaceAfter=8
                        ))
                    else:
                        # Regular text - add to current paragraph
                        current_paragraph.append(line_stripped)
                
                # Flush remaining paragraph
                if current_paragraph:
                    para_text = ' '.join(current_paragraph)
                    para_text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', para_text)
                    para_text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', para_text)
                    elements.append(Paragraph(para_text, styles['BodyText']))
            else:
                elements.append(Paragraph("No content available.", styles['BodyText']))

            # Add metadata section at end
            elements.append(Spacer(1, 20))
            elements.append(HRFlowable(
                width="100%",
                thickness=1,
                color=PDF_COLORS["light"],
                spaceBefore=8,
                spaceAfter=8
            ))
            
            # Document metadata
            meta_items = []
            if generated_at:
                meta_items.append(f"Generated: {generated_at.strftime('%B %d, %Y at %I:%M %p UTC')}")
            if version:
                meta_items.append(f"Version: {version}")
            meta_items.append(f"Card: {card_name}")
            
            for item in meta_items:
                elements.append(Paragraph(item, styles['SmallText']))

            # Build PDF using professional builder with header/footer
            builder = ProfessionalPDFBuilder(
                filename=pdf_path,
                title=brief_title,
                include_logo=True,
                include_ai_disclosure=True
            )
            builder.build(elements)

            logger.info(f"Generated professional brief PDF: {brief_title}")
            return pdf_path

        except Exception as e:
            logger.error(f"Error generating professional brief PDF: {e}")
            raise

    async def generate_brief_pptx(
        self,
        brief_title: str,
        card_name: str,
        executive_summary: str,
        content_markdown: str,
        generated_at: Optional[datetime] = None,
        version: Optional[int] = None
    ) -> str:
        """
        Generate a PowerPoint presentation for an executive brief.

        Args:
            brief_title: Title for the brief (usually card name)
            card_name: Name of the card this brief is for
            executive_summary: Executive summary text
            content_markdown: Full brief content in markdown format
            generated_at: When the brief was generated
            version: Version number of the brief

        Returns:
            Path to the generated PowerPoint file

        Raises:
            Exception: If PowerPoint generation fails
        """
        try:
            logger.info(f"Generating brief PowerPoint: {brief_title}")

            # Create presentation
            prs = Presentation()
            prs.slide_width = PPTX_SLIDE_WIDTH
            prs.slide_height = PPTX_SLIDE_HEIGHT

            # 1. Title slide
            subtitle = f"Executive Brief for {card_name}"
            if version and version > 1:
                subtitle += f" • Version {version}"
            if generated_at:
                subtitle += f" • {generated_at.strftime('%B %d, %Y')}"

            self._add_title_slide(
                prs,
                title=brief_title,
                subtitle=subtitle
            )

            # 2. Executive Summary slide
            if executive_summary:
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)

                # Title bar
                title_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0),
                    PPTX_SLIDE_WIDTH, Inches(1.2)
                )
                title_bar.fill.solid()
                title_bar.fill.fore_color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["secondary"])
                title_bar.line.fill.background()

                # Title
                title_box = slide.shapes.add_textbox(
                    PPTX_MARGIN, Inches(0.3),
                    PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(0.8)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Executive Summary"
                title_para.font.size = Pt(32)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(255, 255, 255)

                # Summary content
                summary_box = slide.shapes.add_textbox(
                    PPTX_MARGIN, Inches(1.6),
                    PPTX_SLIDE_WIDTH - (2 * PPTX_MARGIN), Inches(5.5)
                )
                summary_frame = summary_box.text_frame
                summary_frame.word_wrap = True
                summary_para = summary_frame.paragraphs[0]
                # Truncate if too long for slide
                summary_text = executive_summary[:1500] if len(executive_summary) > 1500 else executive_summary
                summary_para.text = summary_text
                summary_para.font.size = PPTX_BODY_FONT_SIZE
                summary_para.font.color.rgb = self._hex_to_rgb(FORESIGHT_COLORS["dark"])
                summary_para.line_spacing = 1.4

            # 3. Content slides - parse markdown into sections
            if content_markdown:
                sections = self._parse_markdown_sections(content_markdown)

                for section_title, section_content in sections:
                    self._add_description_slide(
                        prs,
                        title=section_title or "Brief Content",
                        description=section_content
                    )

            # Save presentation to temp file
            temp_file = tempfile.NamedTemporaryFile(
                suffix='.pptx',
                delete=False,
                prefix='foresight_brief_'
            )
            prs.save(temp_file.name)

            logger.info(f"Brief PowerPoint generated successfully: {temp_file.name}")
            return temp_file.name

        except Exception as e:
            logger.error(f"Error generating brief PowerPoint: {e}")
            raise

    def _parse_markdown_sections(
        self,
        content_markdown: str,
        max_sections: int = 10
    ) -> List[Tuple[str, str]]:
        """
        Parse markdown content into sections based on headers.

        Args:
            content_markdown: Markdown content to parse
            max_sections: Maximum number of sections to return

        Returns:
            List of (title, content) tuples
        """
        import re

        sections = []
        current_title = "Overview"
        current_content = []

        lines = content_markdown.split('\n')

        for line in lines:
            # Check for headers
            header_match = re.match(r'^(#{1,3})\s+(.+)$', line)
            if header_match:
                # Save previous section if it has content
                if current_content:
                    content_text = '\n'.join(current_content).strip()
                    if content_text:
                        sections.append((current_title, content_text))

                current_title = header_match.group(2)
                current_content = []
            else:
                current_content.append(line)

        # Don't forget the last section
        if current_content:
            content_text = '\n'.join(current_content).strip()
            if content_text:
                sections.append((current_title, content_text))

        # If no sections were found, return all content as one section
        if not sections and content_markdown.strip():
            sections = [("Brief Content", content_markdown.strip())]

        return sections[:max_sections]
